"""Generate fake test documents in a nested hierarchy for the aws-connect E2E test."""
import os, shutil, zipfile
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import openpyxl

STAGE = os.environ.get("STAGE", os.path.join(os.environ["TEMP"], "aws-connect-testdocs"))

# Author stamped into Office document properties (docProps/core.xml). The ingest pipeline reads
# this with stdlib zipfile and writes it to the AI Search 'author' field, so the E2E can assert it.
AUTHOR = "E2E Test Author"


def marker(ext):
    """Distinctive token embedded in each generated file's *content* so the E2E can assert that
    real, correctly-extracted text reached AI Search for that file type. Document Intelligence
    pulls PDF text from pages[].lines[] but Office/HTML text from paragraphs/content, so we verify
    extraction *quality* per extension, not merely that some chunk exists."""
    return f"CONTENTMARKER-{ext.upper()}"


LOREM = (
    "The quarterly results reflect steady growth across all business units. "
    "Revenue increased compared to the prior period driven by strong demand. "
    "Operating margins improved as the team executed on cost discipline. "
    "Management remains optimistic about the outlook for the coming year. "
)

def multipage_pdf(path, title, pages):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    c = canvas.Canvas(path, pagesize=letter)
    w, h = letter
    for p in range(1, pages + 1):
        c.setFont("Helvetica-Bold", 16)
        c.drawString(1 * inch, h - 1 * inch, f"{title} — Page {p} of {pages}")
        c.setFont("Helvetica", 11)
        y = h - 1.5 * inch
        if p == 1:
            c.drawString(1 * inch, y, marker("pdf"))
            y -= 0.28 * inch
        # ~24 lines of text so each page has real, extractable content
        for line in range(24):
            c.drawString(1 * inch, y, f"[p{p} L{line+1}] " + LOREM[: 90])
            y -= 0.28 * inch
        c.showPage()
    c.save()
    print("PDF ", path, f"({pages}p)")


def make_docx(path, title, paras, author=AUTHOR):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    d = Document()
    d.add_heading(title, 0)
    d.add_paragraph(marker("docx"))
    for i in range(paras):
        d.add_paragraph(f"Section {i+1}. " + LOREM * 2)
    d.core_properties.author = author
    d.save(path)
    print("DOCX", path)


def make_txt(path, title, lines):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    ext = os.path.splitext(path)[1].lstrip(".").lower() or "txt"
    with open(path, "w", encoding="utf-8") as f:
        f.write(title + "\n")
        f.write(marker(ext) + "\n\n")
        for i in range(lines):
            f.write(f"{i+1}. " + LOREM + "\n")
    print(ext.upper().ljust(4), path)


def make_md(path, title, lines):
    # Markdown is read directly (no Doc Intelligence), same as .txt; make_txt derives the marker
    # from the .md extension.
    make_txt(path, "# " + title, lines)


def make_pptx(path, title, slides=3, author=AUTHOR):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs = Presentation()
    layout = prs.slide_layouts[5]  # title only
    for s in range(1, slides + 1):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"{title} — Slide {s} of {slides}"
        tb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = (marker("pptx") + " ") if s == 1 else ""
        tf.add_paragraph().text = LOREM * 2
    prs.core_properties.author = author
    prs.save(path)
    print("PPTX", path, f"({slides} slides)")


def make_xlsx(path, title, author=AUTHOR):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = title
    ws["A2"] = marker("xlsx")
    for col, hdr in enumerate(["Item", "Quarter", "Revenue", "Notes"], start=1):
        ws.cell(row=4, column=col, value=hdr)
    for r in range(5, 25):
        ws.cell(row=r, column=1, value=f"Item {r-4}")
        ws.cell(row=r, column=2, value=f"Q{(r % 4) + 1}")
        ws.cell(row=r, column=3, value=1000 + r * 7)
        ws.cell(row=r, column=4, value=LOREM[:60])
    wb.properties.creator = author
    wb.save(path)
    print("XLSX", path)


def make_html(path, title):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    ext = os.path.splitext(path)[1].lstrip(".").lower() or "html"
    body = "\n".join(f"<p>{i+1}. {LOREM}</p>" for i in range(8))
    doc = (f"<!DOCTYPE html><html><head><meta charset='utf-8'><title>{title}</title></head>"
           f"<body><h1>{title}</h1><p>{marker(ext)}</p>{body}</body></html>")
    with open(path, "w", encoding="utf-8") as f:
        f.write(doc)
    print(ext.upper().ljust(4), path)

def make_zip_unsupported(path):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with zipfile.ZipFile(path, "w") as z:
        z.writestr("readme.txt", "unsupported archive to exercise the skip path")
    print("ZIP ", path, "(unsupported)")

def main():
    if os.path.exists(STAGE):
        shutil.rmtree(STAGE)
    os.makedirs(STAGE)
    j = lambda *p: os.path.join(STAGE, *p)

    quarters = ["q1", "q2", "q3", "q4"]
    years = ["2023", "2024"]

    # finance/ (has ACL) — nested reports by year/quarter + policies
    for yr in years:
        for q in quarters:
            multipage_pdf(j("finance", "reports", yr, q, f"{q}-{yr}-earnings.pdf"),
                          f"{q.upper()} {yr} Earnings", 3)
            make_txt(j("finance", "reports", yr, q, f"{q}-{yr}-analyst-notes.txt"),
                     f"Analyst Notes {q.upper()} {yr}", 10)
    for name in ["expense-policy", "travel-policy", "procurement-policy",
                 "audit-controls", "revenue-recognition"]:
        make_docx(j("finance", "policies", f"{name}.docx"),
                  name.replace("-", " ").title(), 5)
    for i in range(1, 31):
        multipage_pdf(j("finance", "invoices", "2024", f"invoice-{i:04d}.pdf"),
                      f"Invoice {i:04d}", 2)

    # hr/ (has ACL) — handbook, onboarding, benefits, deeply nested reviews
    multipage_pdf(j("hr", "handbook", "employee-handbook.pdf"), "Employee Handbook", 6)
    make_zip_unsupported(j("hr", "handbook", "handbook-archive.zip"))  # unsupported -> skip
    for i in range(1, 11):
        multipage_pdf(j("hr", "onboarding", f"welcome-guide-{i:02d}.pdf"),
                      f"New Hire Welcome Guide {i:02d}", 2)
    for name in ["health-plan", "dental-plan", "retirement-401k", "pto-policy",
                 "parental-leave", "wellness-program"]:
        make_docx(j("hr", "benefits", f"{name}.docx"),
                  name.replace("-", " ").title(), 4)
    for dept in ["sales", "support", "operations"]:
        for i in range(1, 7):
            make_txt(j("hr", "reviews", "2024", dept, f"review-{dept}-{i:02d}.txt"),
                     f"Performance Review {dept} {i:02d}", 8)

    # engineering/ (NO ACL entry -> no_acl skip unless bypass), deeply nested
    for area in ["platform", "data", "security"]:
        for i in range(1, 5):
            multipage_pdf(j("engineering", "specs", area, "v2", "internal", "draft",
                            f"{area}-spec-{i:02d}.pdf"),
                          f"{area.title()} Spec {i:02d}", 3)

    print("\nSTAGE:", STAGE)
    total = sum(len(fs) for _, _, fs in os.walk(STAGE))
    print("total files:", total)

if __name__ == "__main__":
    main()
